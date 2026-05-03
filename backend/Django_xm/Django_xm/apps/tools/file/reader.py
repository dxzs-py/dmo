"""
文件读取与解析工具
支持读取用户上传的文件内容，用于基于文件内容回答问题
"""

import os
import base64
from pathlib import Path
from typing import Optional, List, Tuple, Dict, Any

from langchain_core.tools import tool
from langchain_core.documents import Document

import logging

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h",
    ".css", ".scss", ".less", ".yaml", ".yml", ".toml",
    ".ini", ".cfg", ".conf", ".sh", ".bat", ".sql",
    ".go", ".rs", ".rb", ".php", ".swift", ".kt",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}

MAX_FILE_SIZE = 10 * 1024 * 1024
MAX_CONTENT_LENGTH = 50000

IMAGE_MIME_MAP = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".svg": "image/svg+xml",
}


def _read_text_file(file_path: Path) -> str:
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法解码文件: {file_path}")


def _read_pdf_file(file_path: Path) -> str:
    try:
        from langchain_community.document_loaders import PyPDFLoader
        loader = PyPDFLoader(str(file_path))
        pages = loader.load()
        return "\n\n".join(page.page_content for page in pages)
    except ImportError:
        try:
            import fitz
            doc = fitz.open(str(file_path))
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            doc.close()
            return "\n\n".join(text_parts)
        except ImportError:
            raise ValueError("PDF 解析需要安装 PyPDF2 或 PyMuPDF: pip install pypdf 或 pip install PyMuPDF")


def _read_docx_file(file_path: Path) -> str:
    try:
        from langchain_community.document_loaders import Docx2txtLoader
        loader = Docx2txtLoader(str(file_path))
        docs = loader.load()
        return "\n".join(doc.page_content for doc in docs)
    except ImportError:
        try:
            import docx
            doc = docx.Document(str(file_path))
            return "\n".join(para.text for para in doc.paragraphs)
        except ImportError:
            raise ValueError("DOCX 解析需要安装 python-docx: pip install python-docx")


def _read_excel_file(file_path: Path) -> str:
    try:
        import pandas as pd
        ext = file_path.suffix.lower()
        if ext == ".csv":
            df = pd.read_csv(str(file_path))
        else:
            df = pd.read_excel(str(file_path))
        return df.to_string(index=False)
    except ImportError:
        raise ValueError("Excel 解析需要安装 pandas 和 openpyxl: pip install pandas openpyxl")


def _read_pptx_file(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- 幻灯片 {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        return "\n\n".join(text_parts)
    except ImportError:
        raise ValueError("PPTX 解析需要安装 python-pptx: pip install python-pptx")


def _read_image_file(file_path: Path) -> str:
    try:
        from Django_xm.apps.ai_engine.services.llm_factory import get_chat_model

        with open(file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")

        ext = file_path.suffix.lower()
        mime_type = IMAGE_MIME_MAP.get(ext, "image/png")

        model = get_chat_model(model_name="gpt-4o", temperature=0.0)

        from langchain_core.messages import HumanMessage
        message = HumanMessage(content=[
            {"type": "text", "text": "请详细描述这张图片的内容。"},
            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_data}"}},
        ])

        response = model.invoke([message])
        return getattr(response, "content", "无法识别图片内容")
    except Exception as e:
        raise ValueError(f"图片识别失败: {str(e)}")


def is_image_file(file_path: str) -> bool:
    return Path(file_path).suffix.lower() in IMAGE_EXTENSIONS


def is_image_extension(ext: str) -> bool:
    return f".{ext.lstrip('.')}" in IMAGE_EXTENSIONS


def read_file_content(file_path: str) -> str:
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    if not path.is_file():
        raise ValueError(f"路径不是文件: {file_path}")

    if path.stat().st_size > MAX_FILE_SIZE:
        raise ValueError(f"文件过大（超过 {MAX_FILE_SIZE // (1024*1024)}MB）: {file_path}")

    ext = path.suffix.lower()

    if ext in SUPPORTED_EXTENSIONS:
        content = _read_text_file(path)
    elif ext == ".pdf":
        content = _read_pdf_file(path)
    elif ext in (".doc", ".docx"):
        content = _read_docx_file(path)
    elif ext in (".xls", ".xlsx"):
        content = _read_excel_file(path)
    elif ext in (".ppt", ".pptx"):
        content = _read_pptx_file(path)
    elif ext in IMAGE_EXTENSIONS:
        content = _read_image_file(path)
    else:
        try:
            content = _read_text_file(path)
        except Exception:
            raise ValueError(f"不支持的文件类型: {ext}")

    if len(content) > MAX_CONTENT_LENGTH:
        content = content[:MAX_CONTENT_LENGTH] + f"\n\n... [文件内容过长，已截断，原始长度: {len(content)} 字符]"

    return content


def read_file_as_documents(file_path: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Document]:
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    if not path.is_file():
        raise ValueError(f"路径不是文件: {file_path}")

    if path.stat().st_size > MAX_FILE_SIZE:
        raise ValueError(f"文件过大（超过 {MAX_FILE_SIZE // (1024*1024)}MB）: {file_path}")

    ext = path.suffix.lower()
    file_name = path.name

    if ext in IMAGE_EXTENSIONS:
        text_content = _read_image_file(path)
        return [Document(page_content=text_content, metadata={"source": file_name, "file_type": "image"})]

    if ext in SUPPORTED_EXTENSIONS:
        raw_text = _read_text_file(path)
    elif ext == ".pdf":
        try:
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(str(path))
            docs = loader.load()
            for doc in docs:
                doc.metadata["source"] = file_name
            return docs
        except ImportError:
            raw_text = _read_pdf_file(path)
    elif ext in (".doc", ".docx"):
        raw_text = _read_docx_file(path)
    elif ext in (".xls", ".xlsx"):
        raw_text = _read_excel_file(path)
    elif ext in (".ppt", ".pptx"):
        raw_text = _read_pptx_file(path)
    else:
        try:
            raw_text = _read_text_file(path)
        except Exception:
            raise ValueError(f"不支持的文件类型: {ext}")

    if len(raw_text) <= chunk_size:
        return [Document(page_content=raw_text, metadata={"source": file_name, "file_type": ext.lstrip(".")})]

    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""],
        )
        chunks = splitter.create_documents(
            texts=[raw_text],
            metadatas=[{"source": file_name, "file_type": ext.lstrip(".")}],
        )
        logger.info(f"文件 {file_name} 分割为 {len(chunks)} 个片段")
        return chunks
    except ImportError:
        logger.warning("langchain_text_splitters 不可用，返回完整文档")
        return [Document(page_content=raw_text, metadata={"source": file_name, "file_type": ext.lstrip(".")})]


def get_attachment_info(attachment_id: int) -> Dict[str, Any]:
    try:
        from Django_xm.apps.chat.models import ChatAttachment
        attachment = ChatAttachment.objects.select_related('session').get(id=attachment_id)
    except Exception as e:
        raise ValueError(f"找不到附件 (id={attachment_id}): {str(e)}")

    file_path = attachment.file.path
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"附件文件不存在: {attachment.original_name}")

    ext = Path(file_path).suffix.lower()

    return {
        "id": attachment.id,
        "original_name": attachment.original_name,
        "file_type": attachment.file_type,
        "mime_type": attachment.mime_type,
        "file_size": attachment.file_size,
        "file_path": file_path,
        "is_image": ext in IMAGE_EXTENSIONS,
        "ext": ext,
    }


def read_attachment_as_base64(attachment_id: int) -> Tuple[str, str]:
    info = get_attachment_info(attachment_id)
    file_path = info["file_path"]
    ext = info["ext"]
    mime_type = IMAGE_MIME_MAP.get(ext, "image/png")

    with open(file_path, "rb") as f:
        image_data = base64.b64encode(f.read()).decode("utf-8")

    return image_data, mime_type


def read_attachment_as_documents(attachment_id: int, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Document]:
    info = get_attachment_info(attachment_id)
    file_path = info["file_path"]

    docs = read_file_as_documents(file_path, chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    for doc in docs:
        doc.metadata["attachment_id"] = attachment_id
        doc.metadata["original_name"] = info["original_name"]

    return docs


def read_uploaded_attachment(attachment_id: int) -> str:
    try:
        from Django_xm.apps.chat.models import ChatAttachment
        attachment = ChatAttachment.objects.select_related('session').get(id=attachment_id)
    except Exception as e:
        raise ValueError(f"找不到附件 (id={attachment_id}): {str(e)}")

    file_path = attachment.file.path
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"附件文件不存在: {attachment.original_name}")

    content = read_file_content(file_path)

    header = (
        f"📄 文件名: {attachment.original_name}\n"
        f"📁 文件类型: {attachment.file_type}\n"
        f"📏 文件大小: {attachment.file_size} 字节\n"
        f"{'─' * 40}\n\n"
    )

    return header + content


def read_multiple_attachments(attachment_ids: List[int]) -> str:
    parts = []
    for att_id in attachment_ids:
        try:
            content = read_uploaded_attachment(att_id)
            parts.append(content)
        except Exception as e:
            parts.append(f"❌ 读取附件 (id={att_id}) 失败: {str(e)}")

    return "\n\n" + ("=" * 60 + "\n\n").join(parts)


@tool
def file_reader(file_path: str) -> str:
    """
    读取指定路径的文件内容

    支持多种文件格式：文本文件(.txt/.md/.py/.js等)、PDF、Word(.docx)、Excel(.xlsx)、PPT(.pptx)、图片等。
    当用户要求阅读、分析或总结某个文件时使用此工具。

    Args:
        file_path: 文件的完整路径

    Returns:
        文件内容文本

    Example:
        >>> file_reader("/path/to/document.pdf")
        '文件内容...'
    """
    logger.info(f"📄 读取文件: {file_path}")

    try:
        content = read_file_content(file_path)
        logger.info(f"📄 文件读取成功: {len(content)} 字符")
        return content
    except FileNotFoundError:
        return f"错误：文件不存在: {file_path}"
    except ValueError as e:
        return f"错误：{str(e)}"
    except Exception as e:
        error_msg = f"读取文件失败: {str(e)}"
        logger.error(error_msg)
        return error_msg


@tool
def attachment_reader(attachment_id: int) -> str:
    """
    读取用户上传的聊天附件内容

    当用户上传了文件并基于文件内容提问时，使用此工具读取附件内容。
    attachment_id 可以从前端消息的 attachments 字段获取。

    Args:
        attachment_id: 附件的 ID

    Returns:
        附件文件内容

    Example:
        >>> attachment_reader(1)
        '文件内容...'
    """
    logger.info(f"📎 读取附件: id={attachment_id}")

    try:
        content = read_uploaded_attachment(attachment_id)
        logger.info(f"📎 附件读取成功: {len(content)} 字符")
        return content
    except Exception as e:
        error_msg = f"读取附件失败: {str(e)}"
        logger.error(error_msg)
        return error_msg


def get_file_reader_tools():
    return [file_reader, attachment_reader]


FILE_READER_TOOLS = get_file_reader_tools()
